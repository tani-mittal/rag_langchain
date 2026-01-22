from pptx import Presentation
from typing import List
from pathlib import Path
from .base_processor import BaseDocumentProcessor, DocumentChunk

class PPTXProcessor(BaseDocumentProcessor):
    def extract_content(self, file_path: Path) -> List[DocumentChunk]:
        chunks = []
        
        try:
            self.logger.info(f"Processing PPTX: {file_path}")
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = ""
                slide_images = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                    
                    # Extract images
                    try:
                        if hasattr(shape, 'image') and shape.image:
                            image_data = shape.image.blob
                            slide_images.append(image_data)
                    except Exception as e:
                        self.logger.warning(f"Error extracting image from slide {slide_num}: {e}")
                        continue
                
                slide_text = self.clean_text(slide_text)
                
                if slide_text:
                    # Process slide text
                    text_chunks = self.split_text(slide_text)
                    
                    for i, chunk_text in enumerate(text_chunks):
                        chunk_id = self.create_chunk_id(str(file_path), len(chunks))
                        
                        chunk = DocumentChunk(
                            content=chunk_text,
                            metadata={
                                "source": str(file_path),
                                "slide": slide_num + 1,
                                "chunk_index": i,
                                "file_type": "pptx",
                                "total_slides": len(prs.slides),
                                "file_name": file_path.name
                            },
                            chunk_id=chunk_id,
                            source_file=str(file_path),
                            page_number=slide_num + 1,
                            images=self.process_images(slide_images) if slide_images and i == 0 else None
                        )
                        chunks.append(chunk)
            
            self.logger.info(f"Successfully processed PPTX: {len(chunks)} chunks created")
            
        except Exception as e:
            self.logger.error(f"Error processing PPTX {file_path}: {str(e)}")
            raise
        
        return chunks